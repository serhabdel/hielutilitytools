from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel,
    QPushButton, QLineEdit, QFileDialog, QProgressBar,
    QMessageBox, QRadioButton, QButtonGroup, QFrame
)
from PySide6.QtCore import Qt
import os
import fitz
from pptx import Presentation
from pptx.util import Inches
import comtypes.client
from utils.app_theme import AppTheme

class PPTConverter(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.theme = AppTheme()
        self.setup_ui()

    def setup_ui(self):
        layout = QVBoxLayout(self)
        layout.setSpacing(20)
        layout.setContentsMargins(30, 30, 30, 30)

        # Title
        title = QLabel("PPT/PDF Converter")
        title.setStyleSheet("font-size: 24px; color: white; margin-bottom: 20px;")
        layout.addWidget(title)

        # Content Frame
        content_frame = QFrame()
        content_frame.setStyleSheet("""
            QFrame {
                background: #2d2d2d;
                border-radius: 8px;
                padding: 20px;
            }
        """)
        content_layout = QVBoxLayout(content_frame)
        content_layout.setSpacing(15)

        # Conversion Type
        type_label = QLabel("Conversion Type")
        type_label.setStyleSheet("color: #888; font-size: 12px;")
        content_layout.addWidget(type_label)

        radio_layout = QHBoxLayout()
        self.button_group = QButtonGroup(self)
        self.ppt_to_pdf = QRadioButton("PPT to PDF")
        self.pdf_to_ppt = QRadioButton("PDF to PPT")
        radio_style = """
            QRadioButton {
                color: white;
                spacing: 8px;
            }
            QRadioButton::indicator {
                width: 18px;
                height: 18px;
                border: 2px solid #4a90e2;
                border-radius: 9px;
                background: #2d2d2d;
            }
            QRadioButton::indicator:checked {
                background: qradialgradient(cx:0.5, cy:0.5, radius:0.4, fx:0.5, fy:0.5, 
                    stop:0 #4a90e2, stop:1 #2d2d2d);
            }
        """
        self.ppt_to_pdf.setStyleSheet(radio_style)
        self.pdf_to_ppt.setStyleSheet(radio_style)
        self.button_group.addButton(self.ppt_to_pdf)
        self.button_group.addButton(self.pdf_to_ppt)
        self.ppt_to_pdf.setChecked(True)
        radio_layout.addWidget(self.ppt_to_pdf)
        radio_layout.addWidget(self.pdf_to_ppt)
        radio_layout.addStretch()
        content_layout.addLayout(radio_layout)

        # Input File
        input_label = QLabel("Input File")
        input_label.setStyleSheet("color: #888; font-size: 12px;")
        content_layout.addWidget(input_label)

        input_layout = QHBoxLayout()
        self.input_path = QLineEdit()
        self.input_path.setPlaceholderText("Choose input file...")
        self.input_path.setStyleSheet("""
            QLineEdit {
                padding: 8px;
                background: #2d2d2d;
                border: 1px solid #3d3d3d;
                border-radius: 4px;
                color: white;
            }
            QLineEdit:focus {
                border: 1px solid #4a90e2;
            }
        """)
        browse_btn = QPushButton("Browse")
        browse_btn.setStyleSheet("""
            QPushButton {
                padding: 8px 16px;
                background: #4a90e2;
                border: none;
                border-radius: 4px;
                color: white;
            }
            QPushButton:hover {
                background: #357abd;
            }
        """)
        browse_btn.clicked.connect(self.browse_input)
        input_layout.addWidget(self.input_path)
        input_layout.addWidget(browse_btn)
        content_layout.addLayout(input_layout)

        # Output File
        output_label = QLabel("Output File")
        output_label.setStyleSheet("color: #888; font-size: 12px;")
        content_layout.addWidget(output_label)

        output_layout = QHBoxLayout()
        self.output_path = QLineEdit()
        self.output_path.setPlaceholderText("Choose output location...")
        self.output_path.setStyleSheet(self.input_path.styleSheet())
        browse_output_btn = QPushButton("Browse")
        browse_output_btn.setStyleSheet(browse_btn.styleSheet())
        browse_output_btn.clicked.connect(self.browse_output)
        output_layout.addWidget(self.output_path)
        output_layout.addWidget(browse_output_btn)
        content_layout.addLayout(output_layout)

        # Progress Bar
        self.progress = QProgressBar()
        self.progress.setStyleSheet("""
            QProgressBar {
                border: 1px solid #3d3d3d;
                border-radius: 4px;
                text-align: center;
                background: #363636;
            }
            QProgressBar::chunk {
                background-color: #4a90e2;
            }
        """)
        self.progress.hide()
        content_layout.addWidget(self.progress)

        # Convert Button
        convert_btn = QPushButton("Convert")
        convert_btn.setStyleSheet("""
            QPushButton {
                padding: 12px 24px;
                background: #4a90e2;
                border: none;
                border-radius: 4px;
                color: white;
                font-size: 14px;
                font-weight: bold;
            }
            QPushButton:hover {
                background: #357abd;
            }
        """)
        convert_btn.clicked.connect(self.convert)
        content_layout.addWidget(convert_btn)

        layout.addWidget(content_frame)
        layout.addStretch()

    def update_file_paths(self):
        # Clear paths when conversion type changes
        self.input_path.clear()
        self.output_path.clear()

    def browse_input(self):
        file_filter = "PowerPoint files (*.pptx *.ppt)" if self.ppt_to_pdf.isChecked() else "PDF files (*.pdf)"
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Select Input File",
            "",
            file_filter
        )
        if file_path:
            self.input_path.setText(file_path)
            if not self.output_path.text():
                base_name = os.path.splitext(file_path)[0]
                ext = ".pdf" if self.ppt_to_pdf.isChecked() else ".pptx"
                self.output_path.setText(base_name + ext)

    def browse_output(self):
        file_filter = "PDF files (*.pdf)" if self.ppt_to_pdf.isChecked() else "PowerPoint files (*.pptx)"
        file_path, _ = QFileDialog.getSaveFileName(
            self,
            "Save Output File",
            "",
            file_filter
        )
        if file_path:
            self.output_path.setText(file_path)

    def convert_ppt_to_pdf(self, input_path, output_path):
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = True

            ppt = powerpoint.Presentations.Open(input_path)
            ppt.SaveAs(output_path, 32)  # 32 is the PDF format code
            ppt.Close()
            powerpoint.Quit()

            return True
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to convert PPT to PDF: {str(e)}")
            return False

    def convert_pdf_to_ppt(self, input_path, output_path):
        try:
            # Create a new presentation
            prs = Presentation()

            # Open PDF
            pdf_document = fitz.open(input_path)
            total_pages = len(pdf_document)

            for page_num in range(total_pages):
                # Add a slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

                # Convert PDF page to image
                page = pdf_document[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scale for better quality

                # Save page as temporary image
                temp_img = f"temp_slide_{page_num}.png"
                pix.save(temp_img)

                # Add image to slide
                slide.shapes.add_picture(temp_img, 0, 0, prs.slide_width, prs.slide_height)

                # Clean up temp file
                os.remove(temp_img)

                self.progress.setValue(int((page_num + 1) / total_pages * 100))

            # Save presentation
            prs.save(output_path)
            pdf_document.close()

            return True
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to convert PDF to PPT: {str(e)}")
            return False

    def convert(self):
        input_path = self.input_path.text()
        output_path = self.output_path.text()

        if not input_path or not output_path:
            QMessageBox.warning(self, "Error", "Please select both input and output files")
            return

        try:
            self.progress.show()
            self.progress.setValue(0)

            if self.ppt_to_pdf.isChecked():
                success = self.convert_ppt_to_pdf(os.path.abspath(input_path), os.path.abspath(output_path))
            else:
                success = self.convert_pdf_to_ppt(input_path, output_path)

            self.progress.setValue(100)
            self.progress.hide()

            if success:
                QMessageBox.information(
                    self,
                    "Success",
                    f"Conversion completed successfully!\nOutput file: {output_path}"
                )
        except Exception as e:
            self.progress.hide()
            QMessageBox.critical(self, "Error", f"An error occurred: {str(e)}")

